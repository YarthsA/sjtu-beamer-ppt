#!/usr/bin/env python3
"""
pdf_to_pptx.py — Convert Beamer PDF to PPTX
Auto-detects best available engine: WPS (editable text) > PyMuPDF (image fallback).

Usage:
  python pdf_to_pptx.py --input slides.pdf                    # auto-detect engine
  python pdf_to_pptx.py --input slides.pdf --engine wps       # force WPS
  python pdf_to_pptx.py --input slides.pdf --engine pymupdf   # force PyMuPDF image
"""

import argparse
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


def find_wps():
    """Dynamically locate WPS kwpsconvert.exe on this system."""
    # 1. Check PATH
    path = shutil.which("kwpsconvert")
    if path:
        return path

    # 2. Search common install locations via environment variables
    env_candidates = []
    for var in ("LOCALAPPDATA", "PROGRAMFILES", "PROGRAMFILES(X86)"):
        base = os.environ.get(var, "")
        if base:
            env_candidates.append(Path(base) / "kingsoft")

    # 3. Search registry (Windows only)
    try:
        import winreg
        for hive in (winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE):
            for subkey in [
                r"SOFTWARE\Kingsoft\Office",
                r"SOFTWARE\WOW6432Node\Kingsoft\Office",
            ]:
                try:
                    with winreg.OpenKey(hive, subkey) as key:
                        for i in range(winreg.QueryInfoKey(key)[0]):
                            name = winreg.EnumKey(key, i)
                            try:
                                with winreg.OpenKey(key, name) as sub:
                                    val, _ = winreg.QueryValueEx(sub, "InstallPath")
                                    if val:
                                        env_candidates.append(Path(val))
                            except OSError:
                                pass
                except OSError:
                    pass
    except ImportError:
        pass  # not Windows

    # 4. Scan candidate directories for kwpsconvert.exe
    for base in env_candidates:
        if not base.exists():
            continue
        # WPS version dirs are typically under office6/
        for pattern in ["**/office6/kwpsconvert.exe", "**/kwpsconvert.exe"]:
            for match in base.glob(pattern):
                return str(match)

    return None


def convert_wps(pdf_path, output_path, engine="ai"):
    """Convert PDF to editable PPTX using WPS (requires WPS VIP)."""
    wps = find_wps()
    if not wps:
        return False, "WPS not found on this system"

    cmd = [
        wps, "pdf2ppt",
        str(pdf_path),
        "-o", str(output_path),
        "--engine", engine,
        "--image-quality", "best",
    ]

    print(f"  Using WPS ({engine} mode) ...")
    result = subprocess.run(cmd, capture_output=True, timeout=300)
    stdout = result.stdout.decode("utf-8", errors="replace") if result.stdout else ""
    stderr = result.stderr.decode("utf-8", errors="replace") if result.stderr else ""

    if result.returncode == 0 and output_path.exists():
        return True, None
    elif result.returncode in (100, 101):
        return False, "WPS requires VIP account (exit 100/101)"
    else:
        msg = stdout.strip() or stderr.strip() or f"exit code {result.returncode}"
        return False, msg


def convert_pymupdf(pdf_path, output_path, dpi=300):
    """Convert PDF to PPTX as high-res images (universal fallback, not editable)."""
    try:
        import fitz
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        return False, f"Missing dependency: {e}. Run: pip install PyMuPDF python-pptx Pillow"

    doc = fitz.open(str(pdf_path))
    page_count = doc.page_count
    if page_count == 0:
        doc.close()
        return False, "PDF has no pages"

    SLIDE_W_EMU = 12192000  # 13.333 inches = 33.867cm
    SLIDE_H_EMU = 6858000   # 7.5 inches = 19.05cm
    SLIDE_W_IN = 13.333
    SLIDE_H_IN = 7.5

    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]

    slide_w_px = int(dpi * SLIDE_W_IN)
    slide_h_px = int(dpi * SLIDE_H_IN)

    print(f"  Using PyMuPDF ({dpi} DPI, image mode) ...")

    with tempfile.TemporaryDirectory() as tmpdir:
        for i in range(page_count):
            page = doc[i]
            print(f"  Rendering page {i + 1}/{page_count} ...", end="\r")

            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            pix = None

            iw, ih = img.size
            ia, sa = iw / ih, slide_w_px / slide_h_px
            if ia > sa:
                nw, nh = int(iw * slide_h_px / ih), slide_h_px
            else:
                nw, nh = slide_w_px, int(ih * slide_w_px / iw)
            img_r = img.resize((nw, nh), Image.LANCZOS)
            img.close()
            l, t = (nw - slide_w_px) // 2, (nh - slide_h_px) // 2
            img_c = img_r.crop((l, t, l + slide_w_px, t + slide_h_px))
            img_r.close()

            tmp_png = Path(tmpdir) / f"s_{i + 1:04d}.png"
            img_c.save(str(tmp_png), "PNG", optimize=True)
            img_c.close()

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(tmp_png), Inches(0), Inches(0),
                                     Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))

    doc.close()
    prs.save(str(output_path))
    return True, None


def main():
    parser = argparse.ArgumentParser(description="Convert Beamer PDF to PPTX")
    parser.add_argument("--input", "-i", required=True, help="Input PDF path")
    parser.add_argument("--output", "-o", default=None, help="Output PPTX path")
    parser.add_argument("--engine", "-e", default="auto",
                        choices=["auto", "wps", "pymupdf"],
                        help="Conversion engine (default: auto)")
    parser.add_argument("--dpi", "-d", type=int, default=300,
                        help="PyMuPDF render DPI (pymupdf engine only, default: 300)")
    args = parser.parse_args()

    pdf_path = Path(args.input)
    if not pdf_path.exists():
        print(f"Error: PDF not found: {pdf_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output) if args.output else pdf_path.with_suffix(".pptx")

    print(f"Converting: {pdf_path.name}")

    success = False
    used_engine = ""

    # Try WPS (auto or explicit)
    if args.engine in ("auto", "wps"):
        for eng in ["ai", "basic"]:
            ok, err = convert_wps(pdf_path, output_path, engine=eng)
            if ok:
                success = True
                used_engine = f"WPS ({eng})"
                break
            if args.engine == "wps":
                print(f"  WPS failed: {err}", file=sys.stderr)
                sys.exit(1)

    # Fallback to PyMuPDF
    if not success and args.engine in ("auto", "pymupdf"):
        ok, err = convert_pymupdf(pdf_path, output_path, dpi=args.dpi)
        if ok:
            success = True
            used_engine = "PyMuPDF (image)"
        else:
            print(f"  PyMuPDF failed: {err}", file=sys.stderr)
            sys.exit(1)

    if not success:
        print("Error: All conversion engines failed", file=sys.stderr)
        sys.exit(1)

    size_mb = output_path.stat().st_size / (1024 * 1024)
    print(f"\nDone: {output_path}")
    print(f"  Engine: {used_engine}, Size: {size_mb:.1f} MB")


if __name__ == "__main__":
    main()
